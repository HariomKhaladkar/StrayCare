import sys
import collections
import collections.abc
from pptx import Presentation

pptx_path = r"c:\Users\Hariom\Downloads\Nexora2026_IDEA_Presentation_Format.pptx"

try:
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text.replace("\n", " ").strip()
            if text:
                print(f"  Shape {shape.shape_id}: {text[:100]}")
except Exception as e:
    print(f"Error: {e}")
