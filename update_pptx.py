import collections
import collections.abc
from pptx import Presentation

def replace_in_shape(shape, replacements):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        full_text = "".join(run.text for run in paragraph.runs)
        for key, val in replacements.items():
            if key in full_text:
                full_text = full_text.replace(key, val)
                # clear all runs except first
                if len(paragraph.runs) > 0:
                    paragraph.runs[0].text = full_text
                    for i in range(1, len(paragraph.runs)):
                        paragraph.runs[i].text = ""

input_pptx = r"c:\Users\Hariom\Downloads\Nexora2026_IDEA_Presentation_Format.pptx"
output_pptx = r"c:\Users\Hariom\Downloads\Nexora2026_IDEA_Presentation_Format (1).pptx"

try:
    prs = Presentation(input_pptx)
    
    # Define replacements per slide or globally.
    # Since text might be tricky with line breaks, let's just do targeted replacements if the text matches exactly,
    # or just replace the entire text if we identify the shape.

    def update_slide(slide_index, replacements):
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in paragraph.runs)
                for key, val in replacements.items():
                    if key in full_text:
                        full_text = full_text.replace(key, val)
                        if len(paragraph.runs) > 0:
                            paragraph.runs[0].text = full_text
                            for i in range(1, len(paragraph.runs)):
                                paragraph.runs[i].text = ""

    # Slide 1 Replacements
    s1_repl = {
        "Problem Statement Title-": "Problem Statement Title- Stray Animal Rescue & Management",
        "Theme-": "Theme- Civic Tech / Animal Welfare",
        "PS Category-": "PS Category- Software",
        "Team Name (Registered on portal)-": "Team Name (Registered on portal)- Team StrayCare"
    }
    update_slide(0, s1_repl)

    # Slide 2 Replacements
    s2_repl = {
        "IDEA TITLE": "STRAYCARE - AI-Augmented Animal Rescue Platform",
        "Detailed explanation of the proposed solution": "StrayCare is a full-stack, AI-augmented civic-tech platform bridging citizens, NGOs, and administrators. It automates triage using NLP, dispatches rescue resources via geo-spatial Haversine scoring, and provides a Zone Red Hotspot map via K-Means clustering.",
        "How it addresses the problem": "It transforms fragmented rescue operations into a unified, efficient network. Critical cases are auto-triaged, and NGOs are dispatched based on proximity and caseload.",
        "Innovation and uniqueness of the solution": "Runs AI in-browser (zero inference cost), implements an Uber-style NGO dispatch system, and predicts hotspots without external paid APIs."
    }
    update_slide(1, s2_repl)

    # Slide 3 Replacements
    s3_repl = {
        "Technologies to be used (e.g. programming languages, frameworks, hardware)": "Frontend: React.js, Tailwind, Capacitor.js (Android)\nBackend: FastAPI (Python), SQLite\nAI/ML: TensorFlow.js (MobileNet V2), Groq API LLama 3",
        "Methodology and process for implementation (Flow Charts/Images/ working prototype)": "1. Citizen reports case; In-browser AI detects animal species.\n2. Backend runs NLP to sort cases by severity.\n3. Smart Dispatch assigns nearest NGO using Haversine.\n4. Admin monitors analytics and Zone Red hotspots."
    }
    update_slide(2, s3_repl)

    # Slide 4 Replacements
    s4_repl = {
        "Analysis of the feasibility of the idea": "Highly achievable. Zero recurring API costs for core intelligence (runs locally in-browser or pure Python). High reach via Web & Android.",
        "Potential challenges and risks": "Challenges: Poor internet connectivity in remote areas, dependance on active citizen reporting, NGO coordination.",
        "Strategies  for overcoming these challenges": "Strategies: PWA design with offline AI caching; localized awareness drives; seamless Razorpay UPI payments to sustain NGOs financially."
    }
    update_slide(3, s4_repl)

    # Slide 5 Replacements
    s5_repl = {
        "Potential impact on the target audience": "Citizens: Easy platform to report & adopt. NGOs: Automated triage & matching reduces workload. Administrators: Live hotspot maps.",
        "Benefits of the solution (social, economic, environmental, etc.)": "Social: Fosters community compassion. Economic: Saves costs through optimized resource dispatch. Environmental: Sustainable stray population management and disease control."
    }
    update_slide(4, s5_repl)

    # Slide 6 Replacements
    s6_repl = {
        "Details / Links of the reference and research work": "- [1] Stray Animal Mobile App: Lacks multi-tier tracking\n- [3] StandForPaw: Community abuse reporting without backend\n- [5] FurRescue: AI breed detection but focused only on lost/found\nStrayCare scales this by integrating unified dashboards, K-Means clustering, and Haversine routing."
    }
    update_slide(5, s6_repl)

    prs.save(output_pptx)
    print("Presentation successfully updated!")

except Exception as e:
    import traceback
    traceback.print_exc()

